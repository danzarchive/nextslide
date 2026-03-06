#!/usr/bin/env python3
"""
NextSlide — PPTX Content Extractor

Extracts text, structure, and images from an existing .pptx file
so the AI can convert it into a NextSlide HTML presentation.

Usage:
    python extract-pptx.py <path-to-pptx> [--images-dir <dir>] [--format json|text]

Requirements:
    pip install python-pptx Pillow

Output:
    Structured text or JSON with slide content, layout hints, and image paths.
"""

import argparse
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx is not installed.", file=sys.stderr)
    print("Install it with: pip install python-pptx Pillow", file=sys.stderr)
    sys.exit(1)


def emu_to_inches(emu):
    """Convert EMU (English Metric Units) to inches."""
    if emu is None:
        return 0
    return round(emu / 914400, 3)


def emu_to_pct(emu, total_emu):
    """Convert EMU to percentage of slide dimension."""
    if emu is None or total_emu is None or total_emu == 0:
        return "0%"
    return f"{round((emu / total_emu) * 100, 1)}%"


def pt_to_float(pt_val):
    """Convert Pt value to float points."""
    if pt_val is None:
        return None
    return round(pt_val.pt, 1)


def extract_alignment(alignment):
    """Convert PP_ALIGN enum to string."""
    if alignment is None:
        return "left"
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "left",
    }
    return mapping.get(alignment, "left")


def rgb_to_hex(rgb_color):
    """Convert RGBColor to hex string."""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color}"
    except Exception:
        return None


def guess_layout(slide):
    """Guess the NextSlide layout type from slide content."""
    shapes = slide.shapes
    has_image = any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        or (hasattr(shape, "image") and shape.image is not None)
        for shape in shapes
    )
    text_shapes = [s for s in shapes if s.has_text_frame]
    text_count = len(text_shapes)

    # Check for placeholder types
    has_title = any(
        hasattr(s, "placeholder_format")
        and s.placeholder_format is not None
        and s.placeholder_format.idx == 0
        for s in shapes
    )

    # Simple heuristics
    if text_count <= 2 and has_title and not has_image:
        total_text = " ".join(
            s.text_frame.text for s in text_shapes if s.has_text_frame
        )
        if len(total_text) < 100:
            return "title"
    if has_image and text_count > 0:
        return "image-left"
    if text_count == 1 and has_title:
        return "section"
    if text_count > 3:
        return "bullets"
    return "content"


def extract_shape(
    shape, slide_width_emu, slide_height_emu, images_dir=None, slide_num=0, shape_idx=0
):
    """Extract content from a single shape."""
    result = {
        "type": "unknown",
        "position": {
            "x": emu_to_pct(shape.left, slide_width_emu),
            "y": emu_to_pct(shape.top, slide_height_emu),
            "w": emu_to_pct(shape.width, slide_width_emu),
            "h": emu_to_pct(shape.height, slide_height_emu),
        },
    }

    # Image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
        hasattr(shape, "image") and shape.image is not None
    ):
        result["type"] = "image"
        if images_dir and hasattr(shape, "image") and shape.image is not None:
            ext = shape.image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_filename = f"slide{slide_num}_img{shape_idx}.{ext}"
            img_path = os.path.join(images_dir, img_filename)
            with open(img_path, "wb") as f:
                f.write(shape.image.blob)
            result["src"] = img_path
            result["alt"] = shape.name or f"Image from slide {slide_num}"
        else:
            result["src"] = f"[image: {shape.name}]"
            result["alt"] = shape.name or ""
        return result

    # Table
    if shape.has_table:
        result["type"] = "table"
        table = shape.table
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            rows.append(cells)
        result["data"] = rows
        return result

    # Chart
    if shape.has_chart:
        result["type"] = "chart"
        chart = shape.chart
        result["chartType"] = (
            str(chart.chart_type).split(".")[-1].lower()
            if chart.chart_type
            else "unknown"
        )
        chart_data = []
        try:
            for series in chart.series:
                series_data = {
                    "name": str(series.name) if hasattr(series, "name") else "Series",
                    "values": [v for v in series.values]
                    if hasattr(series, "values")
                    else [],
                }
                chart_data.append(series_data)
        except Exception:
            pass
        result["data"] = chart_data
        return result

    # Text
    if shape.has_text_frame:
        result["type"] = "text"
        tf = shape.text_frame
        full_text = tf.text.strip()

        if not full_text:
            return None  # Skip empty text boxes

        result["content"] = full_text

        # Guess role from placeholder or font size
        role = "body"
        if (
            hasattr(shape, "placeholder_format")
            and shape.placeholder_format is not None
        ):
            idx = shape.placeholder_format.idx
            if idx == 0:
                role = "title"
            elif idx == 1:
                role = "subtitle"
            elif idx == 10:
                role = "subtitle"
        result["role"] = role

        # Extract first paragraph style as representative
        if tf.paragraphs:
            para = tf.paragraphs[0]
            style = {}
            if para.font.size:
                style["fontSize"] = round(para.font.size.pt, 0)
            if para.font.name:
                style["fontFace"] = para.font.name
            if para.font.bold:
                style["bold"] = True
            if para.font.italic:
                style["italic"] = True
            if para.font.color and para.font.color.rgb:
                style["color"] = f"#{para.font.color.rgb}"
            if para.alignment:
                style["align"] = extract_alignment(para.alignment)
            if style:
                result["style"] = style

        return result

    # Generic shape (rectangle, circle, etc.)
    if shape.shape_type in (
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.FREEFORM,
    ):
        result["type"] = "shape"
        result["shape"] = "rect"  # Default; hard to determine exact shape
        if hasattr(shape, "fill") and shape.fill.type is not None:
            try:
                result["style"] = {"fill": f"#{shape.fill.fore_color.rgb}"}
            except Exception:
                pass
        return result

    return None


def extract_notes(slide):
    """Extract speaker notes from a slide."""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        if notes_tf:
            text = notes_tf.text.strip()
            if text:
                return text
    return None


def extract_presentation(pptx_path, images_dir=None):
    """Extract all content from a PPTX file."""
    prs = Presentation(pptx_path)

    # Slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    meta = {
        "title": os.path.splitext(os.path.basename(pptx_path))[0],
        "slideWidth": round(slide_width / 914400, 3),
        "slideHeight": round(slide_height / 914400, 3),
        "slideCount": len(prs.slides),
    }

    # Try to extract title from first slide
    if prs.slides:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
                and shape.has_text_frame
            ):
                meta["title"] = shape.text_frame.text.strip()
                break

    # Extract core properties
    if prs.core_properties:
        if prs.core_properties.author:
            meta["author"] = prs.core_properties.author
        if prs.core_properties.created:
            meta["date"] = str(prs.core_properties.created.date())

    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slideNumber": slide_num,
            "layout": guess_layout(slide),
            "elements": [],
        }

        # Extract shapes
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                element = extract_shape(
                    shape, slide_width, slide_height, images_dir, slide_num, shape_idx
                )
                if element:
                    slide_data["elements"].append(element)
            except Exception as e:
                print(f"  ⚠ Slide {slide_num}, shape {shape_idx}: {e}", file=sys.stderr)

        # Extract notes
        notes = extract_notes(slide)
        if notes:
            slide_data["notes"] = notes

        slides.append(slide_data)

    return {"meta": meta, "slides": slides}


def format_as_text(data):
    """Format extracted data as human-readable text for AI consumption."""
    lines = []
    meta = data["meta"]

    lines.append(f"# {meta['title']}")
    lines.append("")
    if meta.get("author"):
        lines.append(f"Author: {meta['author']}")
    if meta.get("date"):
        lines.append(f"Date: {meta['date']}")
    lines.append(f"Slides: {meta['slideCount']}")
    lines.append(f'Dimensions: {meta["slideWidth"]}" × {meta["slideHeight"]}"')
    lines.append("")
    lines.append("---")
    lines.append("")

    for slide in data["slides"]:
        lines.append(f"## Slide {slide['slideNumber']} [{slide['layout']}]")
        lines.append("")

        for el in slide["elements"]:
            if el["type"] == "text":
                role = el.get("role", "body")
                content = el.get("content", "")
                style_parts = []
                if el.get("style"):
                    s = el["style"]
                    if s.get("fontSize"):
                        style_parts.append(f"{s['fontSize']}pt")
                    if s.get("fontFace"):
                        style_parts.append(s["fontFace"])
                    if s.get("bold"):
                        style_parts.append("bold")
                    if s.get("color"):
                        style_parts.append(s["color"])
                style_str = f" ({', '.join(style_parts)})" if style_parts else ""
                lines.append(f"[{role}]{style_str}: {content}")

            elif el["type"] == "image":
                lines.append(
                    f"[image]: {el.get('src', 'unknown')} — {el.get('alt', '')}"
                )

            elif el["type"] == "table":
                lines.append("[table]:")
                for row in el.get("data", []):
                    lines.append(f"  | {' | '.join(row)} |")

            elif el["type"] == "chart":
                lines.append(f"[chart: {el.get('chartType', 'unknown')}]")
                for series in el.get("data", []):
                    lines.append(
                        f"  {series.get('name', 'Series')}: {series.get('values', [])}"
                    )

            elif el["type"] == "shape":
                lines.append(f"[shape: {el.get('shape', 'rect')}]")

        if slide.get("notes"):
            lines.append("")
            lines.append(f"Notes: {slide['notes']}")

        lines.append("")
        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description="Extract content from a .pptx file for NextSlide conversion.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python extract-pptx.py presentation.pptx
  python extract-pptx.py deck.pptx --format json
  python extract-pptx.py deck.pptx --images-dir ./extracted-images
  python extract-pptx.py deck.pptx --format json > content.json
        """,
    )
    parser.add_argument("pptx_path", help="Path to the .pptx file")
    parser.add_argument(
        "--images-dir",
        help="Directory to save extracted images (created if needed)",
        default=None,
    )
    parser.add_argument(
        "--format",
        choices=["text", "json"],
        default="text",
        help="Output format: 'text' (default, human-readable) or 'json' (structured)",
    )
    args = parser.parse_args()

    pptx_path = os.path.abspath(args.pptx_path)
    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    if not pptx_path.lower().endswith(".pptx"):
        print("Warning: File does not have .pptx extension", file=sys.stderr)

    # Create images directory if needed
    images_dir = None
    if args.images_dir:
        images_dir = os.path.abspath(args.images_dir)
        os.makedirs(images_dir, exist_ok=True)

    print(f"\n📄 NextSlide PPTX Extractor", file=sys.stderr)
    print(f"━━━━━━━━━━━━━━━━━━━━━━━━━━", file=sys.stderr)
    print(f"  Input: {pptx_path}", file=sys.stderr)
    if images_dir:
        print(f"  Images: {images_dir}", file=sys.stderr)
    print(f"  Format: {args.format}\n", file=sys.stderr)

    # Extract content
    data = extract_presentation(pptx_path, images_dir)

    print(f"  ✓ Extracted {data['meta']['slideCount']} slides", file=sys.stderr)
    total_elements = sum(len(s["elements"]) for s in data["slides"])
    print(f"  ✓ Found {total_elements} elements", file=sys.stderr)
    if images_dir:
        img_count = sum(
            1 for s in data["slides"] for e in s["elements"] if e["type"] == "image"
        )
        print(f"  ✓ Saved {img_count} images to {images_dir}", file=sys.stderr)
    print("", file=sys.stderr)

    # Output
    if args.format == "json":
        print(json.dumps(data, indent=2, ensure_ascii=False))
    else:
        print(format_as_text(data))


if __name__ == "__main__":
    main()
